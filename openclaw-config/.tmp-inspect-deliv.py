import os, subprocess, json

D = "/root/.openclaw/workspace/deliverables/note_x6lx8zsntb3eqcx5ydp4nkf4"
files = sorted(os.listdir(D))
print(f"=== {len(files)} files in {D} ===\n")
for fname in files:
    path = os.path.join(D, fname)
    size = os.path.getsize(path)
    print(f"--- {fname} ({size} bytes) ---")
    ext = fname.lower().split(".")[-1]
    if ext in ("txt", "md", "csv"):
        with open(path, "rb") as f:
            data = f.read(800)
        print(data.decode("utf-8", errors="replace"))
    elif ext == "pdf":
        # Use pdftotext if available, fall back to strings
        try:
            r = subprocess.run(["pdftotext", "-layout", path, "-"],
                                capture_output=True, timeout=10)
            txt = r.stdout.decode("utf-8", errors="replace")
            print(txt[:1200] if txt else "(pdftotext returned empty)")
        except FileNotFoundError:
            r = subprocess.run(["strings", path], capture_output=True, timeout=10)
            print(r.stdout.decode("utf-8", errors="replace")[:800])
    elif ext == "docx":
        try:
            import zipfile, re
            with zipfile.ZipFile(path) as z:
                with z.open("word/document.xml") as f:
                    xml = f.read().decode("utf-8", errors="replace")
            txt = re.sub(r"<[^>]+>", " ", xml)
            txt = re.sub(r"\s+", " ", txt).strip()
            print(txt[:1200])
        except Exception as e:
            print(f"docx parse err: {e}")
    elif ext == "xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(path, read_only=True, data_only=True)
            for sheet in wb.sheetnames[:3]:
                ws = wb[sheet]
                print(f"  sheet: {sheet}")
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= 8: break
                    print(f"    {row}")
        except Exception as e:
            print(f"xlsx parse err: {e}")
    elif ext == "pptx":
        try:
            from pptx import Presentation
            prs = Presentation(path)
            print(f"  slides: {len(prs.slides)}")
            for i, slide in enumerate(prs.slides[:5]):
                print(f"  slide {i+1}:")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        print(f"    {shape.text[:200]}")
        except Exception as e:
            print(f"pptx parse err: {e}")
    print()
